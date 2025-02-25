from pptx import Presentation
from pptx.util import Inches

# 创建 PPTX 文件
prs = Presentation()

# **封面页**
slide_layout = prs.slide_layouts[0]  # 标题幻灯片
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AI-Based Resume Matching System"
subtitle.text = "Enhancing Job Matching with AI & NLP\nPresenter: Zhenghua (Sean) Mu"

# **Problem Statement**
slide = prs.slides.add_slide(prs.slide_layouts[5])  # 仅标题
title = slide.shapes.title
title.text = "Problem Statement"

text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
text_frame = text_box.text_frame
text_frame.text = ("🔴 International students struggle to find jobs\n"
                   "🔴 Recruiters spend too much time screening resumes\n"
                   "✅ AI-powered system extracts resume details\n"
                   "✅ Matches candidates with jobs using AI similarity search")

# **System Architecture**
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "System Architecture & Tech Stack"

text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
text_frame = text_box.text_frame
text_frame.text = ("🛠 Frontend: React + Redux + TypeScript + TailwindCSS\n"
                   "🛠 Backend: FastAPI + SQLite\n"
                   "🛠 AI Model: all-mpnet-base-v2 (Sentence-BERT)\n"
                   "🛠 Vector Search: FAISS for job similarity matching")

# **Data Processing**
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Data Processing & Resume Parsing"

text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
text_frame = text_box.text_frame
text_frame.text = ("✅ Using Regex for email, skills extraction\n"
                   "✅ Future Plan: Enhance with spaCy NER for better accuracy\n"
                   "✅ Store parsed data in SQLite for efficient querying")

# **Job Matching Process**
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Job Matching with FAISS & NLP"

text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
text_frame = text_box.text_frame
text_frame.text = ("🔍 Convert Resume & Job Descriptions into Vectors\n"
                   "🔍 Store Job Vectors in FAISS Index\n"
                   "🔍 Use Cosine Similarity to Find Best Matches\n"
                   "📊 Results: Ranked job recommendations with similarity scores")

# **Challenges & Solutions**
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Challenges & Solutions"

text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
text_frame = text_box.text_frame
text_frame.text = ("⚠ CORS Issues → Fixed with FastAPI Middleware\n"
                   "⚠ File Upload Issues → Switched to Form-Data\n"
                   "⚠ FAISS Vector Mismatch → Fixed embedding dimensions\n"
                   "⚠ Resume Parsing Accuracy → Future enhancement with spaCy")

# **Future Enhancements**
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Future Enhancements & Next Steps"

text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
text_frame = text_box.text_frame
text_frame.text = ("✅ Improve Resume Parsing with spaCy NER\n"
                   "✅ Expand FAISS Index to include soft skills\n"
                   "✅ Build Web Dashboard for Recruiters\n"
                   "✅ Deploy on AWS/GCP for scalability")

# **Summary & Q&A**
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Summary & Q&A"

text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
text_frame = text_box.text_frame
text_frame.text = ("🏆 We built an AI-driven job matching system\n"
                   "🏆 Successfully solved resume parsing & job retrieval\n"
                   "🏆 Matching is real-time (~milliseconds per query)\n"
                   "💬 Any questions? Happy to discuss!")

# **保存 PPTX 文件**
pptx_path = "AI_Resume_Matching_Presentation.pptx"
prs.save(pptx_path)

print(f"PPTX 文件已生成: {pptx_path}")
